from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from services.slide_schema import Slide
import re
import requests
from io import BytesIO
from core.logger import get_logger

logger = get_logger("academic_layouts")

class EquationLayout:
    """Layout specifically designed for mathematical equations"""
    
    def render(self, slide_data: Slide, pptx_slide):
        # Set up title with mathematical styling
        title_shape = pptx_slide.shapes.title
        title_shape.text = slide_data.title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)  # Larger title
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(21, 67, 96)  # Dark blue
        
        # Add a subtle underline/divider
        divider = pptx_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(2), Inches(1.2), Inches(6), Inches(0.03)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = RGBColor(41, 128, 185)  # Medium blue
        divider.line.fill.background()  # No outline
        
        # Create content area with enhanced styling for mathematical equations
        if len(pptx_slide.placeholders) > 1:
            content_shape = pptx_slide.placeholders[1]
        else:
            # If no placeholder found, create a custom text box with better positioning
            content_shape = pptx_slide.shapes.add_textbox(
                Inches(1), Inches(1.8), Inches(8), Inches(4)
            )
            
        content_shape.left = Inches(1)
        content_shape.width = Inches(8)
        
        tf = content_shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Pt(12)   # Add some left margin
        tf.margin_right = Pt(12)  # Add some right margin
        
        # Process content looking for equation markers
        content_text = []
        if slide_data.content:
            for point in slide_data.content:
                if hasattr(point, 'text'):
                    content_text.append(point.text)
                else:
                    content_text.append(str(point))
        elif slide_data.bullets:
            content_text = slide_data.bullets
        
        # Add each paragraph with enhanced formatting for equations
        for text in content_text:
            p = tf.add_paragraph()
            
            # Apply enhanced equation formatting
            formatted_text = self._format_equation_text(text)
            p.text = formatted_text
            
            # Apply different styling based on content type
            if self._is_equation_line(text):
                # Equations get special treatment
                p.alignment = PP_ALIGN.CENTER
                p.font.bold = True
                p.font.size = Pt(24)
                p.font.color.rgb = RGBColor(39, 174, 96)  # Green for equations
                p.space_before = Pt(12)  # Add space before equations
                p.space_after = Pt(12)   # Add space after equations
            else:
                # Regular text
                p.alignment = PP_ALIGN.LEFT
                p.font.size = Pt(20)
                p.level = 0
                # Add bullet style
                p.bullet.enabled = True
        
        # Add image if available (e.g. graph or plot)
        if slide_data.image_url:
            try:
                response = requests.get(slide_data.image_url, timeout=10)
                response.raise_for_status()
                
                img_stream = BytesIO(response.content)
                pptx_slide.shapes.add_picture(
                    img_stream, 
                    Inches(5), 
                    Inches(3.5), 
                    width=Inches(4)
                )
            except Exception as e:
                logger.error(f"Failed to add equation image: {e}")
        else:
            # Check if we should add a visual matrix representation
            if "matrix" in slide_data.title.lower() or any("matrix" in bullet.lower() for bullet in content_text if isinstance(bullet, str)):
                # Add a visual matrix representation
                self._add_matrix_visualization(pptx_slide)
    
    def _format_equation_text(self, text):
        """Format text with equation markers for better readability"""
        # In a full implementation, this would convert LaTeX to Office equation format
        # For now, we'll enhance the formatting for better display
        formatted = text
        
        # Advanced LaTeX-like substitutions with more mathematical symbols
        substitutions = [
            # Greek letters
            (r'\alpha', 'α'), (r'\beta', 'β'), (r'\gamma', 'γ'), (r'\delta', 'δ'),
            (r'\epsilon', 'ε'), (r'\zeta', 'ζ'), (r'\eta', 'η'), (r'\theta', 'θ'),
            (r'\iota', 'ι'), (r'\kappa', 'κ'), (r'\lambda', 'λ'), (r'\mu', 'μ'),
            (r'\nu', 'ν'), (r'\xi', 'ξ'), (r'\pi', 'π'), (r'\rho', 'ρ'),
            (r'\sigma', 'σ'), (r'\tau', 'τ'), (r'\upsilon', 'υ'), (r'\phi', 'φ'),
            (r'\chi', 'χ'), (r'\psi', 'ψ'), (r'\omega', 'ω'),
            
            # Capital Greek letters
            (r'\Gamma', 'Γ'), (r'\Delta', 'Δ'), (r'\Theta', 'Θ'), (r'\Lambda', 'Λ'),
            (r'\Xi', 'Ξ'), (r'\Pi', 'Π'), (r'\Sigma', 'Σ'), (r'\Phi', 'Φ'),
            (r'\Psi', 'Ψ'), (r'\Omega', 'Ω'),
            
            # Operators and symbols
            (r'\times', '×'), (r'\div', '÷'), (r'\pm', '±'), (r'\mp', '∓'),
            (r'\leq', '≤'), (r'\geq', '≥'), (r'\neq', '≠'), (r'\approx', '≈'),
            (r'\equiv', '≡'), (r'\sum', '∑'), (r'\prod', '∏'), (r'\int', '∫'),
            (r'\partial', '∂'), (r'\infty', '∞'), (r'\nabla', '∇'), (r'\forall', '∀'),
            (r'\exists', '∃'), (r'\in', '∈'), (r'\subset', '⊂'), (r'\supset', '⊃'),
            (r'\cup', '∪'), (r'\cap', '∩'), (r'\emptyset', '∅'), (r'\sqrt', '√'),
            
            # Matrix-specific formatting
            (r'_{ij}', '₍ᵢⱼ₎'), (r'_{i}', '₍ᵢ₎'), (r'_{j}', '₍ⱼ₎'), (r'_{k}', '₍ₖ₎'),
            (r'\cdot', '·'),
            
            # Clean up LaTeX markers
            (r'\left', ''), (r'\right', ''), 
            
            # Convert simple subscripts and superscripts
            ('_0', '₀'), ('_1', '₁'), ('_2', '₂'), ('_3', '₃'), ('_4', '₄'),
            ('_5', '₅'), ('_6', '₆'), ('_7', '₇'), ('_8', '₈'), ('_9', '₉'),
            ('^0', '⁰'), ('^1', '¹'), ('^2', '²'), ('^3', '³'), ('^4', '⁴'),
            ('^5', '⁵'), ('^6', '⁶'), ('^7', '⁷'), ('^8', '⁸'), ('^9', '⁹')
        ]
        
        # Apply substitutions
        for old, new in substitutions:
            formatted = formatted.replace(old, new)
        
        # Clean up any remaining LaTeX dollar signs
        formatted = formatted.replace('$', '')
        
        # Handle matrix notation specifically for better readability
        if '\\sum_{k}' in text or '\\sum_' in text:
            formatted = formatted.replace('\\sum_{k}', '∑ₖ ')
            formatted = formatted.replace('\\sum_', '∑')
        
        return formatted
        
    def _add_matrix_visualization(self, slide):
        """Add a visual representation of matrix operations to enhance the slide"""
        # Matrix A - Create a rectangle with cells
        self._create_matrix_shape(slide, "A", Inches(0.9), Inches(4.8), rows=2, cols=3, 
                                 cell_size=Inches(0.6), color=RGBColor(41, 128, 185))
        
        # Matrix B - Create a rectangle with cells
        self._create_matrix_shape(slide, "B", Inches(2.9), Inches(4.8), rows=2, cols=2, 
                                 cell_size=Inches(0.6), color=RGBColor(52, 152, 219))
        
        # Matrix Result (AB) - Create a rectangle with cells
        self._create_matrix_shape(slide, "AB", Inches(5.4), Inches(4.8), rows=2, cols=2, 
                                 cell_size=Inches(0.6), color=RGBColor(39, 174, 96))
        
        # Add operation symbol
        times_symbol = slide.shapes.add_textbox(
            Inches(2.4), Inches(5.1), Inches(0.4), Inches(0.4)
        )
        times_para = times_symbol.text_frame.paragraphs[0]
        times_para.text = "×"
        times_para.font.size = Pt(28)
        times_para.font.bold = True
        times_para.alignment = PP_ALIGN.CENTER
        
        # Add equals symbol
        equals_symbol = slide.shapes.add_textbox(
            Inches(4.9), Inches(5.1), Inches(0.4), Inches(0.4)
        )
        equals_para = equals_symbol.text_frame.paragraphs[0]
        equals_para.text = "="
        equals_para.font.size = Pt(28)
        equals_para.font.bold = True
        equals_para.alignment = PP_ALIGN.CENTER
        
        # Add caption explaining the visual
        caption = slide.shapes.add_textbox(
            Inches(0.9), Inches(6.3), Inches(6.5), Inches(0.5)
        )
        caption_text = caption.text_frame.paragraphs[0]
        caption_text.text = "Visual representation of matrix multiplication: dimensions (2×3) × (3×2) = (2×2)"
        caption_text.font.italic = True
        caption_text.font.size = Pt(14)
        caption_text.alignment = PP_ALIGN.CENTER
    
    def _create_matrix_shape(self, slide, label, left, top, rows, cols, cell_size, color):
        """Create a visual matrix with the given dimensions"""
        # Add matrix label
        label_shape = slide.shapes.add_textbox(
            left, top - Inches(0.4), Inches(0.6), Inches(0.3)
        )
        label_text = label_shape.text_frame.paragraphs[0]
        label_text.text = label
        label_text.font.size = Pt(20)
        label_text.font.bold = True
        label_text.alignment = PP_ALIGN.CENTER
        
        # Matrix container (brackets)
        # Left bracket
        left_bracket = slide.shapes.add_textbox(
            left - Inches(0.2), top - Inches(0.1), Inches(0.2), Inches(rows * cell_size.inches + 0.2)
        )
        lb_text = left_bracket.text_frame.paragraphs[0]
        lb_text.text = "["
        lb_text.font.size = Pt(44)
        lb_text.font.bold = True
        left_bracket.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Right bracket
        right_bracket = slide.shapes.add_textbox(
            left + (cols * cell_size.inches), top - Inches(0.1), 
            Inches(0.2), Inches(rows * cell_size.inches + 0.2)
        )
        rb_text = right_bracket.text_frame.paragraphs[0]
        rb_text.text = "]"
        rb_text.font.size = Pt(44)
        rb_text.font.bold = True
        right_bracket.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Create grid cells
        for r in range(rows):
            for c in range(cols):
                cell = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left + (c * cell_size.inches), top + (r * cell_size.inches),
                    cell_size, cell_size
                )
                cell.fill.solid()
                cell.fill.fore_color.rgb = color
                cell.line.color.rgb = RGBColor(255, 255, 255)  # White grid lines
                
                # Add cell value (aᵢⱼ, bᵢⱼ, etc.)
                if label == "A":
                    value = f"a{r+1}{c+1}"
                elif label == "B":
                    value = f"b{r+1}{c+1}"
                elif label == "AB":
                    value = f"c{r+1}{c+1}"
                else:
                    value = f"{label}{r+1}{c+1}"
                    
                # Add text to cell
                tf = cell.text_frame
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.text = value
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(255, 255, 255)  # White text
    
    def _is_equation_line(self, text):
        """Check if this line appears to be an equation"""
        # Check for equation indicators
        equation_markers = [
            '=', '+', '-', '×', '÷', '/', '*', '^', 
            '$', '\\frac', '\\sum', '\\int', '\\cdot',
            '_', '{', '}', '\\left', '\\right', '\\prod',
            '\\lim', '\\alpha', '\\beta', '\\gamma', '\\theta',
            '\\lambda', '\\delta', '\\sigma', '\\omega'
        ]
        
        # Count mathematical symbols
        math_symbol_count = sum(text.count(marker) for marker in equation_markers)
        
        # Check for specific patterns that indicate equations
        equation_patterns = [
            r'\$.+\$',                   # LaTeX equation delimiters
            r'=.+[a-zA-Z0-9]',           # Equations with = sign
            r'\(.+\).+\(.+\)',           # Multiple parenthetical expressions
            r'\w+_{[a-zA-Z0-9]+}',       # Subscript notation
            r'\w+\^\{[a-zA-Z0-9]+\}',    # Superscript notation
            r'\\\w+\{.+\}',              # LaTeX command with arguments
            r'\s\\sum|\s\\prod|\s\\int', # Math operators with space before
            r'\)\s*=|\}\s*=',            # Right parenthesis or brace followed by =
        ]
        
        # Check if any equation patterns match
        has_equation_pattern = any(re.search(pattern, text) for pattern in equation_patterns)
        
        # If we have significant math symbols or a matching pattern
        return (math_symbol_count >= 2 or 
                has_equation_pattern or 
                (text.strip().startswith('$') and text.strip().endswith('$')) or
                ('\\sum' in text) or ('\\int' in text) or 
                (text.count('_') > 1 and text.count('{') > 0))

class CodeLayout:
    """Layout specifically designed for code snippets"""
    
    def render(self, slide_data: Slide, pptx_slide):
        # Set up title
        title_shape = pptx_slide.shapes.title
        title_shape.text = slide_data.title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        
        # Create a shape to hold the code snippet
        code_shape = pptx_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), 
            Inches(1.8),
            Inches(8),
            Inches(4)
        )
        
        # Set shape properties for code display
        code_shape.fill.solid()
        code_shape.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray background
        code_shape.line.color.rgb = RGBColor(200, 200, 200)
        
        # Add text frame for code
        code_frame = code_shape.text_frame
        code_frame.word_wrap = True
        code_frame.margin_left = Pt(12)
        code_frame.margin_right = Pt(12)
        code_frame.margin_top = Pt(12)
        code_frame.margin_bottom = Pt(12)
        
        # Extract code from content
        code_text = self._extract_code(slide_data)
        
        # Add code with monospace font
        p = code_frame.paragraphs[0]
        p.text = code_text
        p.font.name = "Courier New"  # Monospace font
        p.font.size = Pt(14)
        
        # Add explanation text if available
        if slide_data.notes:
            notes_shape = pptx_slide.shapes.add_textbox(
                Inches(1),
                Inches(6),
                Inches(8),
                Inches(1)
            )
            notes_frame = notes_shape.text_frame
            notes_frame.word_wrap = True
            
            p = notes_frame.add_paragraph()
            p.text = slide_data.notes
            p.font.italic = True
            p.font.size = Pt(14)
    
    def _extract_code(self, slide_data: Slide) -> str:
        """Extract code snippet from slide data"""
        # Try to find code blocks in content
        if slide_data.content:
            for point in slide_data.content:
                if hasattr(point, 'text'):
                    text = point.text
                    # Check for code block markers
                    if text.strip().startswith('```') or text.strip().startswith('def ') or text.strip().startswith('class '):
                        return text.replace('```', '').strip()
            
            # If no explicit code block, join all content
            return '\n'.join(str(point) for point in slide_data.content)
        
        elif slide_data.bullets:
            # Join bullet points as code lines
            return '\n'.join(slide_data.bullets)
        
        return "# Code example not found"

class TaxonomyLayout:
    """Layout specifically designed for biological classification/taxonomy"""
    
    def render(self, slide_data: Slide, pptx_slide):
        # Set up title
        title_shape = pptx_slide.shapes.title
        title_shape.text = slide_data.title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        
        # Create a hierarchical classification diagram
        taxonomy_levels = [
            "Kingdom", "Phylum", "Class", "Order", "Family", "Genus", "Species"
        ]
        
        # Extract taxonomy data from content
        taxonomy_data = self._extract_taxonomy_data(slide_data, taxonomy_levels)
        
        # Create the hierarchical diagram
        self._create_taxonomy_diagram(pptx_slide, taxonomy_data, taxonomy_levels)
    
    def _extract_taxonomy_data(self, slide_data: Slide, taxonomy_levels):
        """Extract taxonomy information from slide content"""
        taxonomy_data = {}
        
        # Process content or bullets
        content_text = []
        if slide_data.content:
            for point in slide_data.content:
                if hasattr(point, 'text'):
                    content_text.append(point.text)
                else:
                    content_text.append(str(point))
        elif slide_data.bullets:
            content_text = slide_data.bullets
        
        # Look for taxonomy level markers
        for text in content_text:
            for level in taxonomy_levels:
                if level.lower() in text.lower():
                    # Extract the value after the level name
                    parts = text.split(":", 1)
                    if len(parts) > 1:
                        taxonomy_data[level] = parts[1].strip()
                    else:
                        # Try splitting by level name
                        parts = text.lower().split(level.lower(), 1)
                        if len(parts) > 1:
                            taxonomy_data[level] = parts[1].strip()
        
        return taxonomy_data
    
    def _create_taxonomy_diagram(self, slide, taxonomy_data, taxonomy_levels):
        """Create a visual taxonomy diagram"""
        # Diagram parameters
        start_top = Inches(1.8)
        width = Inches(3)
        height = Inches(0.6)
        vertical_gap = Inches(0.8)
        
        # Colors for gradient effect
        colors = [
            RGBColor(220, 230, 242),  # Light blue
            RGBColor(198, 224, 180),  # Light green
            RGBColor(255, 230, 153),  # Light yellow
            RGBColor(248, 203, 173),  # Light orange
            RGBColor(230, 184, 175),  # Light red
            RGBColor(204, 192, 218),  # Light purple
            RGBColor(180, 199, 231),  # Light blue again
        ]
        
        # Create shapes for each taxonomy level
        for i, level in enumerate(taxonomy_levels):
            if level in taxonomy_data:
                # Calculate position (centered horizontally)
                left = Inches(4.5) - (width / 2)
                top = start_top + (i * vertical_gap)
                
                # Create shape
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left, top, width, height
                )
                
                # Set colors
                shape.fill.solid()
                shape.fill.fore_color.rgb = colors[i % len(colors)]
                shape.line.color.rgb = RGBColor(0, 0, 0)
                
                # Add text
                tf = shape.text_frame
                tf.text = f"{level}: {taxonomy_data[level]}"
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                tf.paragraphs[0].font.size = Pt(16)
                tf.paragraphs[0].font.bold = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                # Connect with previous level (except first)
                if i > 0:
                    connector = slide.shapes.add_connector(
                        MSO_SHAPE.LINE,
                        left + (width / 2),
                        top,
                        left + (width / 2),
                        top - Inches(0.2)
                    )
                    connector.line.color.rgb = RGBColor(0, 0, 0)

# Add additional academic layouts for text dense, text sparse, etc.
class TextDenseLayout:
    """Layout for text-heavy slides"""
    
    def render(self, slide_data: Slide, pptx_slide):
        title_shape = pptx_slide.shapes.title
        title_shape.text = slide_data.title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        
        # Use smaller font and tighter spacing for dense text
        content_shape = None
        for shape in pptx_slide.shapes:
            if shape.has_text_frame and shape != title_shape:
                content_shape = shape
                break
                
        if not content_shape:
            content_shape = pptx_slide.shapes.add_textbox(
                Inches(0.5), Inches(1.8), Inches(9), Inches(5)
            )
            
        tf = content_shape.text_frame
        tf.clear()
        tf.word_wrap = True
        
        # Add content with smaller font size
        if slide_data.bullets:
            for bullet in slide_data.bullets:
                p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(16)  # Smaller font
                p.level = 0
        
        elif slide_data.content:
            for item in slide_data.content:
                p = tf.add_paragraph()
                if hasattr(item, 'text'):
                    p.text = item.text
                    p.level = item.level if hasattr(item, 'level') else 0
                    
                    # Add sub-points if available
                    if hasattr(item, 'sub_points') and item.sub_points:
                        for sub in item.sub_points:
                            sub_p = tf.add_paragraph()
                            sub_p.text = sub
                            sub_p.level = p.level + 1
                            sub_p.font.size = Pt(14)  # Even smaller for sub-points
                else:
                    p.text = str(item)
                    p.level = 0
                
                p.font.size = Pt(16)  # Smaller font for all main points

class TextSparseLayout:
    """Layout for slides with minimal text"""
    
    def render(self, slide_data: Slide, pptx_slide):
        # Make title more prominent
        title_shape = pptx_slide.shapes.title
        title_shape.text = slide_data.title
        title_shape.text_frame.paragraphs[0].font.size = Pt(40)  # Larger title
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        content_shape = None
        for shape in pptx_slide.shapes:
            if shape.has_text_frame and shape != title_shape:
                content_shape = shape
                break
                
        if not content_shape:
            content_shape = pptx_slide.shapes.add_textbox(
                Inches(1.5), Inches(2.5), Inches(7), Inches(3)
            )
            
        tf = content_shape.text_frame
        tf.clear()
        tf.word_wrap = True
        
        # Center the sparse content and make it larger
        if slide_data.bullets:
            for bullet in slide_data.bullets:
                p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(28)  # Larger font
                p.alignment = PP_ALIGN.CENTER
                
        elif slide_data.content:
            for item in slide_data.content:
                p = tf.add_paragraph()
                if hasattr(item, 'text'):
                    p.text = item.text
                else:
                    p.text = str(item)
                p.font.size = Pt(28)  # Larger font
                p.alignment = PP_ALIGN.CENTER
                
class ConclusionLayout:
    """Layout specifically designed for conclusion slides"""
    
    def render(self, slide_data: Slide, pptx_slide):
        # Make title prominent
        title_shape = pptx_slide.shapes.title
        title_shape.text = slide_data.title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Add a decorative line under the title
        line = pptx_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(3), Inches(1.7), Inches(4), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(52, 152, 219)  # Blue line
        line.line.fill.background()  # No outline
        
        # Content in a nice box
        content_box = pptx_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.5), Inches(2.0), Inches(7), Inches(4)
        )
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = RGBColor(245, 245, 245)  # Light gray
        content_box.line.color.rgb = RGBColor(200, 200, 200)  # Gray border
        
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(24)
        tf.margin_right = Pt(24)
        tf.margin_top = Pt(24)
        tf.margin_bottom = Pt(24)
        
        # Add header paragraph
        header_p = tf.add_paragraph()
        header_p.text = "Key Takeaways"
        header_p.font.size = Pt(24)
        header_p.font.bold = True
        header_p.alignment = PP_ALIGN.CENTER
        
        # Add bullet points with spacing between them
        if slide_data.bullets:
            for bullet in slide_data.bullets:
                p = tf.add_paragraph()
                p.text = "• " + bullet
                p.font.size = Pt(20)
                p.space_after = Pt(12)  # Add space between points
                
        elif slide_data.content:
            for item in slide_data.content:
                p = tf.add_paragraph()
                if hasattr(item, 'text'):
                    p.text = "• " + item.text
                else:
                    p.text = "• " + str(item)
                p.font.size = Pt(20)
                p.space_after = Pt(12)  # Add space between points
