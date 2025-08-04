import os
import uuid
import requests
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from services.slide_schema import Deck, Slide, BulletPoint
from services.layout_intelligence import LayoutIntelligence
from services.theme_manager import ThemeManager
from core.logger import get_logger
from io import BytesIO
from PIL import Image

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), '..', 'templates', 'template.pptx')
logger = get_logger("ppt_builder")

# --- Diagram service stub ---
def generate_diagram_image(description: str) -> BytesIO:
    # TODO: Implement with Mermaid, Graphviz, or external API
    # For now, return a blank image
    from PIL import Image
    img = Image.new('RGB', (400, 200), color = (255, 255, 255))
    buf = BytesIO()
    img.save(buf, format='PNG')
    buf.seek(0)
    return buf

# --- Plugin-style layout system ---
class BaseLayout:
    def render(self, slide_data: Slide, pptx_slide):
        raise NotImplementedError

class TitleBulletsLayout(BaseLayout):
    def render(self, slide_data: Slide, pptx_slide):
        title_shape = pptx_slide.shapes.title
        title_shape.text = slide_data.title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        if slide_data.bullets:
            content_shape = None
            for shape in pptx_slide.shapes:
                if shape.has_text_frame and shape != title_shape:
                    content_shape = shape
                    break
            if content_shape:
                tf = content_shape.text_frame
                tf.clear()
                for bullet in slide_data.bullets:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.font.size = Pt(20)
                    p.level = 0
                    p.alignment = PP_ALIGN.LEFT
        if slide_data.notes:
            pptx_slide.notes_slide.notes_text_frame.text = slide_data.notes
        # Images handled by PPTBuilder

class ImageLayout(BaseLayout):
    def render(self, slide_data: Slide, pptx_slide):
        title_shape = pptx_slide.shapes.title
        title_shape.text = slide_data.title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        # Insert first image (if any)
        if slide_data.images:
            img_url = slide_data.images[0]
            try:
                img_bytes = requests.get(img_url, timeout=5).content
                image_stream = BytesIO(img_bytes)
                # Place image in center of slide
                pptx_slide.shapes.add_picture(image_stream, Inches(2), Inches(2), width=Inches(4))
            except Exception as e:
                logger.error(f"Failed to embed image: {img_url} - {e}")
        if slide_data.notes:
            pptx_slide.notes_slide.notes_text_frame.text = slide_data.notes

class DiagramLayout(BaseLayout):
    def render(self, slide_data: Slide, pptx_slide):
        title_shape = pptx_slide.shapes.title
        title_shape.text = slide_data.title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        # Insert first diagram (if any)
        if slide_data.diagrams:
            desc = slide_data.diagrams[0]
            try:
                diagram_img = generate_diagram_image(desc)
                pptx_slide.shapes.add_picture(diagram_img, Inches(2), Inches(2), width=Inches(4))
            except Exception as e:
                logger.error(f"Failed to embed diagram: {desc} - {e}")
        if slide_data.notes:
            pptx_slide.notes_slide.notes_text_frame.text = slide_data.notes

# Import academic layouts
from services.academic_layouts import (
    EquationLayout, CodeLayout, TaxonomyLayout,
    TextDenseLayout, TextSparseLayout, ConclusionLayout
)

# Layout registry
LAYOUT_REGISTRY = {
    # Basic layouts
    "title_bullets": TitleBulletsLayout(),
    "image": ImageLayout(),
    "diagram": DiagramLayout(),
    
    # Academic layouts
    "equation": EquationLayout(),
    "code": CodeLayout(),
    "taxonomy": TaxonomyLayout(),
    "text_dense": TextDenseLayout(),
    "text_sparse": TextSparseLayout(),
    "conclusion": ConclusionLayout()
}

def get_layout(slide_data: Slide):
    if slide_data.type and slide_data.type in LAYOUT_REGISTRY:
        return LAYOUT_REGISTRY[slide_data.type]
    if slide_data.images:
        return LAYOUT_REGISTRY["image"]
    if slide_data.diagrams:
        return LAYOUT_REGISTRY["diagram"]
    return LAYOUT_REGISTRY["title_bullets"]

class PPTBuilder:
    def __init__(self, template_path: str = TEMPLATE_PATH, theme: str = "professional"):
        self.template_path = template_path
        
        # Get theme colors from ThemeManager
        theme_data = ThemeManager.get_theme_colors(theme)
        self.colors = {
            'primary': theme_data['primary'],
            'secondary': theme_data['secondary'],
            'accent': theme_data['accent'],
            'background': theme_data['background'],
            'text': theme_data['text'],
            'light_text': RGBColor(127, 140, 141)  # Light gray (keeping for backward compatibility)
        }
        self.fonts = {
            'title': theme_data['title_font'],
            'body': theme_data['body_font']
        }

    def build(self, deck: Deck, use_template: bool = True, job_id: str = None) -> BytesIO:
        try:
            logger.info(f"Building PPTX. use_template={use_template}, slides={len(deck.slides)}")
            
            # Emit streaming event if job_id provided
            if job_id:
                from services.streaming_service import streaming_service
                streaming_service.emit_event(job_id, "layout_processing", {
                    "message": "Applying intelligent layout selection...",
                    "step": "layout_analysis"
                })
            
            # Apply intelligent layout selection
            layout_engine = LayoutIntelligence()
            deck = layout_engine.determine_optimal_layouts(deck)
            logger.info(f"Applied intelligent layout selection to deck")
            
            # Pre-process slides to handle content overflow
            processed_deck = self._preprocess_slides_for_overflow(deck)
            logger.info(f"Pre-processed slides to handle overflow. Original slides: {len(deck.slides)}, Processed slides: {len(processed_deck.slides)}")
            
            if job_id:
                streaming_service.emit_event(job_id, "slides_processing", {
                    "message": f"Processing {len(processed_deck.slides)} slides...",
                    "total_slides": len(processed_deck.slides),
                    "step": "slide_processing"
                })
            
            if use_template:
                if not os.path.exists(self.template_path):
                    raise FileNotFoundError(f"PPTX template not found: {self.template_path}")
                prs = Presentation(self.template_path)
            else:
                prs = Presentation()
            
            # Clear existing slides if using template
            slide_count = len(prs.slides)
            for i in range(slide_count - 1, -1, -1):
                self._delete_slide(prs, i)
            
            for i, slide_data in enumerate(processed_deck.slides):
                logger.info(f"Processing slide {i+1}: {slide_data.title} (type: {slide_data.type})")
                
                # Emit progress for each slide
                if job_id:
                    streaming_service.emit_event(job_id, "slide_progress", {
                        "message": f"Creating slide {i+1}: {slide_data.title}",
                        "slide_number": i + 1,
                        "slide_title": slide_data.title,
                        "slide_type": slide_data.type,
                        "progress": round((i + 1) / len(processed_deck.slides) * 100, 1),
                        "step": "slide_creation"
                    })
                
                self._create_enhanced_slide(prs, slide_data)
            
            if job_id:
                streaming_service.emit_event(job_id, "finalizing", {
                    "message": "Finalizing presentation...",
                    "step": "finalization"
                })
            
            # Save to BytesIO instead of local file
            pptx_stream = BytesIO()
            prs.save(pptx_stream)
            pptx_stream.seek(0)
            logger.info(f"PPTX generated in memory")
            return pptx_stream
        except Exception as e:
            logger.error(f"Failed to build PPTX: {e}")
            raise
            
    def _preprocess_slides_for_overflow(self, deck: Deck) -> Deck:
        """Pre-process slides to handle content overflow by splitting extremely content-heavy slides"""
        from copy import deepcopy
        
        processed_deck = deepcopy(deck)
        new_slides = []
        
        # Analyze each slide for potential overflow
        for i, slide in enumerate(processed_deck.slides):
            # Skip title slide
            if i == 0 or slide.type == "title":
                new_slides.append(slide)
                continue
                
            # Check if slide has excessive content
            needs_splitting = False
            content_length = 0
            bullet_count = 0
            
            if slide.content:
                bullet_count = len(slide.content) if isinstance(slide.content, list) else 0
                for point in slide.content:
                    if hasattr(point, 'text'):
                        if hasattr(point.text, '__len__'):
                            content_length += len(point.text)
                        else:
                            content_length += len(str(point.text))
                        if hasattr(point, 'sub_points') and point.sub_points and isinstance(point.sub_points, list):
                            bullet_count += len(point.sub_points)
                            for sub in point.sub_points:
                                content_length += len(str(sub))
                    else:
                        content_length += len(str(point))
            elif slide.bullets:
                bullet_count = len(slide.bullets) if isinstance(slide.bullets, list) else 0
                if isinstance(slide.bullets, list):
                    content_length = sum(len(str(bullet)) for bullet in slide.bullets)
                else:
                    content_length = 0
            
            # Determine if slide needs splitting
            needs_splitting = (content_length > 800 or bullet_count > 9)
            
            if needs_splitting and (slide.content or slide.bullets):
                logger.info(f"Splitting content-heavy slide: {slide.title} ({content_length} chars, {bullet_count} bullets)")
                
                # Create split slides
                if slide.content and isinstance(slide.content, list):
                    # Split complex content
                    midpoint = len(slide.content) // 2
                    
                    # First part
                    first_slide = deepcopy(slide)
                    first_slide.content = slide.content[:midpoint]
                    first_slide.title = slide.title
                    new_slides.append(first_slide)
                    
                    # Second part
                    second_slide = deepcopy(slide)
                    second_slide.content = slide.content[midpoint:]
                    second_slide.title = f"{slide.title} (continued)"
                    new_slides.append(second_slide)
                    
                elif slide.bullets and isinstance(slide.bullets, list):
                    # Split simple bullets
                    midpoint = len(slide.bullets) // 2
                    
                    # First part
                    first_slide = deepcopy(slide)
                    first_slide.bullets = slide.bullets[:midpoint]
                    first_slide.title = slide.title
                    new_slides.append(first_slide)
                    
                    # Second part
                    second_slide = deepcopy(slide)
                    second_slide.bullets = slide.bullets[midpoint:]
                    second_slide.title = f"{slide.title} (continued)"
                    new_slides.append(second_slide)
            else:
                # No need to split
                new_slides.append(slide)
        
        # Update deck with new slide sequence
        processed_deck.slides = new_slides
        return processed_deck
    
    def _delete_slide(self, prs, slide_index):
        """Delete slide by index"""
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[slide_index])
    
    def _create_enhanced_slide(self, prs, slide_data: Slide):
        """Create a slide with enhanced formatting"""
        slide_type = slide_data.type or 'content'
        
        if slide_type == 'title':
            self._create_title_slide(prs, slide_data)
        elif slide_type == 'conclusion':
            self._create_conclusion_slide(prs, slide_data)
        else:
            self._create_content_slide(prs, slide_data)
    
    def _create_title_slide(self, prs, slide_data: Slide):
        """Create title slide with enhanced styling"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title formatting
        title = slide.shapes.title
        title.text = slide_data.title
        title_paragraph = title.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(44)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = self.colors['primary']
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Subtitle formatting
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = slide_data.subtitle or slide_data.notes or ""
            if subtitle.text:
                subtitle_paragraph = subtitle.text_frame.paragraphs[0]
                subtitle_paragraph.font.size = Pt(24)
                subtitle_paragraph.font.color.rgb = self.colors['secondary']
                subtitle_paragraph.alignment = PP_ALIGN.CENTER
    
    def _create_content_slide(self, prs, slide_data: Slide):
        """Create content slide with bullet points and formatting"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title formatting
        title = slide.shapes.title
        title.text = slide_data.title
        self._format_title(title)
        
        # Determine layout based on content
        has_image = slide_data.image_url or (slide_data.images and len(slide_data.images) > 0)
        has_diagram = slide_data.diagram_type and slide_data.diagram_data
        
        if has_image and has_diagram:
            # Both image and diagram - use compact layout
            self._create_mixed_content_layout(slide, slide_data)
        elif has_image:
            # Image only - use side-by-side layout
            self._create_image_content_layout(slide, slide_data)
        elif has_diagram:
            # Diagram only - use full width for diagram
            self._create_diagram_content_layout(slide, slide_data)
        else:
            # Text only - use full content area
            self._create_text_only_layout(slide, slide_data)
        
        # Add notes
        if slide_data.notes:
            slide.notes_slide.notes_text_frame.text = slide_data.notes
    
    def _create_text_only_layout(self, slide, slide_data: Slide):
        """Create enhanced layout with adaptive text content and optimal spacing"""
        # Calculate optimal layout based on content characteristics
        text_density = self._calculate_text_density(slide_data)
        
        # Adaptive content area sizing
        if text_density > 0.8:  # Very dense content
            left, width = Inches(0.3), Inches(9.4)
            top, height = Inches(1.3), Inches(5.4)
        elif text_density > 0.5:  # Dense content
            left, width = Inches(0.4), Inches(9.2)
            top, height = Inches(1.4), Inches(5.2)
        else:  # Normal content
            left, width = Inches(0.5), Inches(9.0)
            top, height = Inches(1.5), Inches(5.0)
        
        content_placeholder = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        
        if content_placeholder:
            # Configure content area with adaptive dimensions
            content_placeholder.left = left
            content_placeholder.width = width
            content_placeholder.top = top
            content_placeholder.height = height
            
            # Enhanced content handling based on density
            if text_density > 0.7:
                # Very dense content - use overflow-safe textbox
                self._create_overflow_safe_textbox(slide, slide_data)
            else:
                # Normal content handling with enhanced formatting
                if slide_data.content:
                    self._add_enhanced_bullet_points(content_placeholder, slide_data.content)
                elif slide_data.bullets:
                    self._add_simple_bullet_points(content_placeholder, slide_data.bullets)
                else:
                    # Handle paragraph-style content
                    self._add_paragraph_content(content_placeholder, slide_data)
        else:
            # Create custom text box if no placeholder available
            self._create_overflow_safe_textbox(slide, slide_data)
    
    def _create_image_content_layout(self, slide, slide_data: Slide):
        """Create adaptive layout with image and text side by side with enhanced spacing"""
        # Calculate optimal layout dimensions based on content
        text_density = self._calculate_text_density(slide_data)
        
        # Adaptive text area sizing based on content volume
        if text_density > 0.7:  # High density content
            text_width = Inches(5.2)
            image_width = Inches(3.3)
            image_left = Inches(5.7)
        elif text_density > 0.4:  # Medium density content
            text_width = Inches(4.8)
            image_width = Inches(3.7)
            image_left = Inches(5.3)
        else:  # Low density content
            text_width = Inches(4.3)
            image_width = Inches(4.2)
            image_left = Inches(4.8)
        
        # Create or adjust content area
        content_placeholder = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        
        if content_placeholder:
            # Adaptive content area positioning
            content_placeholder.left = Inches(0.4)
            content_placeholder.width = text_width
            content_placeholder.top = Inches(1.4)
            content_placeholder.height = Inches(5.2)
            
            # Enhanced content handling with adaptive truncation
            if slide_data.content:
                self._add_enhanced_bullet_points(content_placeholder, slide_data.content)
            elif slide_data.bullets:
                self._add_simple_bullet_points(content_placeholder, slide_data.bullets)
        else:
            # Create custom text box if no placeholder available
            txBox = slide.shapes.add_textbox(Inches(0.4), Inches(1.4), text_width, Inches(5.2))
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(18)
            tf.margin_right = Pt(18)
            tf.margin_top = Pt(12)
            tf.margin_bottom = Pt(12)
            
            # Add content with proper formatting
            if slide_data.content:
                self._add_enhanced_bullet_points(txBox, slide_data.content)
            elif slide_data.bullets:
                self._add_simple_bullet_points(txBox, slide_data.bullets)
        
        # Add image with calculated positioning
        image_url = slide_data.image_url or (slide_data.images[0] if slide_data.images else None)
        if image_url:
            self._add_image_to_slide(slide, image_url, position='right', 
                                   custom_pos=(image_left, Inches(1.6), image_width, Inches(4.8)))
    
    def _create_diagram_content_layout(self, slide, slide_data: Slide):
        """Create layout with diagram and minimal text"""
        # Small content area for key points
        content_placeholder = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        
        if content_placeholder:
            # Compact content area at top
            content_placeholder.left = Inches(0.5)
            content_placeholder.width = Inches(9)
            content_placeholder.top = Inches(1.5)
            content_placeholder.height = Inches(2)
            
            # Show only first 3 points to save space
            content = slide_data.content[:3] if slide_data.content else slide_data.bullets[:3]
            if slide_data.content:
                self._add_enhanced_bullet_points(content_placeholder, content)
            else:
                self._add_simple_bullet_points(content_placeholder, content)
        
        # Add diagram below text
        self._add_diagram_to_slide(slide, slide_data)
    
    def _create_mixed_content_layout(self, slide, slide_data: Slide):
        """Create enhanced compact layout with text, image, and diagram"""
        # Calculate optimal layout based on content density
        text_density = self._calculate_text_density(slide_data)
        
        # Adaptive layout dimensions
        if text_density > 0.6:  # High density - prioritize text space
            text_width = Inches(4.2)
            text_height = Inches(2.3)
            image_size = (Inches(2.5), Inches(1.8))
            image_pos = (Inches(6.8), Inches(1.6))
        else:  # Lower density - balanced layout
            text_width = Inches(3.8)
            text_height = Inches(2.0)
            image_size = (Inches(2.8), Inches(2.0))
            image_pos = (Inches(6.5), Inches(1.6))
        
        # Enhanced text area positioning
        content_placeholder = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        
        if content_placeholder:
            # Optimized content area positioning
            content_placeholder.left = Inches(0.4)
            content_placeholder.width = text_width
            content_placeholder.top = Inches(1.4)
            content_placeholder.height = text_height
            
            # Selective content based on density
            if text_density > 0.7:
                content = slide_data.content[:1] if slide_data.content else slide_data.bullets[:2]
            else:
                content = slide_data.content[:2] if slide_data.content else slide_data.bullets[:3]
                
            if slide_data.content:
                self._add_enhanced_bullet_points(content_placeholder, content)
            else:
                self._add_simple_bullet_points(content_placeholder, content)
        
        # Add compact image with calculated positioning
        image_url = slide_data.image_url or (slide_data.images[0] if slide_data.images else None)
        if image_url:
            self._add_compact_image_to_slide(slide, image_url, custom_pos=(*image_pos, *image_size))
        
        # Add compact diagram positioned to avoid overlap
        self._add_compact_diagram_to_slide(slide, slide_data, offset_top=text_height + Inches(0.3))
    
    def _create_conclusion_slide(self, prs, slide_data: Slide):
        """Create conclusion slide"""
        self._create_content_slide(prs, slide_data)  # Same as content slide
    
    def _format_title(self, title_shape):
        """Format slide title with enhanced adaptive sizing and styling"""
        text_frame = title_shape.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.margin_left = Pt(12)
        text_frame.margin_right = Pt(12)
        text_frame.margin_top = Pt(8)
        text_frame.margin_bottom = Pt(8)
        
        # Apply enhanced formatting to paragraph
        title_paragraph = text_frame.paragraphs[0] 
        title_text = title_paragraph.text if title_paragraph.text else ""
        
        # Adaptive font sizing based on title length
        if len(title_text) > 80:
            font_size = 28  # Very long titles
        elif len(title_text) > 50:
            font_size = 32  # Long titles
        elif len(title_text) > 30:
            font_size = 36  # Medium titles  
        else:
            font_size = 40  # Short titles
        
        title_paragraph.font.size = Pt(font_size)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = self.colors['primary']
        title_paragraph.alignment = PP_ALIGN.LEFT
        title_paragraph.line_spacing = 1.1
        title_paragraph.space_after = Pt(6)
    
    def _add_enhanced_bullet_points(self, content_shape, bullet_points):
        """Add formatted bullet points with enhanced structure and text overflow prevention"""
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # Enhanced text frame configuration for better formatting
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.margin_left = Pt(18)
        text_frame.margin_right = Pt(18) 
        text_frame.margin_top = Pt(12)
        text_frame.margin_bottom = Pt(12)
        
        # Calculate adaptive maximum points based on content complexity
        total_text_length = sum(len(str(point.get('text', str(point)) if isinstance(point, dict) else 
                                   getattr(point, 'text', str(point)))) for point in bullet_points)
        
        # Adaptive content limits based on total text volume
        if total_text_length > 2000:
            max_points = 4  # Very dense content
            max_sub_points = 1
        elif total_text_length > 1200:
            max_points = 5  # Dense content
            max_sub_points = 2
        elif total_text_length > 600:
            max_points = 6  # Medium content
            max_sub_points = 3
        else:
            max_points = 7  # Light content
            max_sub_points = 4
        
        for i, point in enumerate(bullet_points[:max_points]):
            if isinstance(point, dict):
                # Handle BulletPoint structure
                main_text = point.get('text', str(point))
                sub_points = point.get('sub_points', [])
                level = point.get('level', 0)
            elif hasattr(point, 'text'):  # BulletPoint object
                main_text = point.text
                sub_points = point.sub_points or []
                level = point.level
            else:
                # Handle simple string
                main_text = str(point)
                sub_points = []
                level = 0
            
            # Adaptive text truncation based on position and level
            if level == 0:
                max_length = 180 if i < 3 else 120  # First 3 points get more space
            else:
                max_length = 100
                
            if len(main_text) > max_length:
                main_text = main_text[:max_length-3] + "..."
                
            # Add main bullet point
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = main_text
            p.level = level
            self._format_bullet_point(p, level)
            
            # Add sub-points with improved formatting
            if sub_points and level < 2:  # Limit nesting depth
                for j, sub_point in enumerate(sub_points[:max_sub_points]):
                    sub_text = str(sub_point)
                    
                    # Progressive truncation for sub-points
                    sub_max_length = max(80 - (j * 10), 40)  # Shorter for later sub-points
                    if len(sub_text) > sub_max_length:
                        sub_text = sub_text[:sub_max_length-3] + "..."
                    
                    sub_p = text_frame.add_paragraph()
                    sub_p.text = sub_text
                    sub_p.level = level + 1
                    self._format_bullet_point(sub_p, level + 1)
                    
        # Add content summary indicator if truncated
        if len(bullet_points) > max_points:
            p = text_frame.add_paragraph()
            remaining = len(bullet_points) - max_points
            p.text = f"... and {remaining} more point{'s' if remaining > 1 else ''}"
            p.level = 0
            p.font.italic = True
            p.font.size = Pt(12)
            p.font.color.rgb = self.colors['light_text']
    
    def _add_simple_bullet_points(self, content_shape, bullets):
        """Add simple bullet points with enhanced auto-fit and adaptive formatting"""
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # Enhanced text frame configuration
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.margin_left = Pt(18)
        text_frame.margin_right = Pt(18)
        text_frame.margin_top = Pt(12)
        text_frame.margin_bottom = Pt(12)
        
        # Calculate total content length for adaptive limits
        total_length = sum(len(str(bullet)) for bullet in bullets)
        
        # Adaptive bullet limits based on total content
        if total_length > 1500:
            max_bullets = 4
            max_length_per_bullet = 120
        elif total_length > 1000:
            max_bullets = 5
            max_length_per_bullet = 150
        elif total_length > 600:
            max_bullets = 6
            max_length_per_bullet = 180
        else:
            max_bullets = 8
            max_length_per_bullet = 220
        
        for i, bullet in enumerate(bullets[:max_bullets]):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # Adaptive text truncation
            bullet_text = str(bullet)
            if len(bullet_text) > max_length_per_bullet:
                bullet_text = bullet_text[:max_length_per_bullet-3] + "..."
                
            p.text = bullet_text
            p.level = 0
            self._format_bullet_point(p, 0)
        
        # Add summary if content was truncated
        if len(bullets) > max_bullets:
            p = text_frame.add_paragraph()
            remaining = len(bullets) - max_bullets
            p.text = f"... plus {remaining} additional point{'s' if remaining > 1 else ''}"
            p.font.italic = True
            p.font.size = Pt(12)
            p.font.color.rgb = self.colors['light_text']
    
    def _add_paragraph_content(self, content_shape, slide_data: Slide):
        """Add paragraph-style content with enhanced formatting for better readability"""
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # Enhanced text frame configuration
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.margin_left = Pt(20)
        text_frame.margin_right = Pt(20)
        text_frame.margin_top = Pt(15)
        text_frame.margin_bottom = Pt(15)
        
        # Handle different content types as paragraphs
        content_text = ""
        if hasattr(slide_data, 'description') and slide_data.description:
            content_text = slide_data.description
        elif slide_data.content and isinstance(slide_data.content, str):
            content_text = slide_data.content
        elif slide_data.bullets and len(slide_data.bullets) == 1:
            content_text = slide_data.bullets[0]
        
        if content_text:
            # Split long content into paragraphs for better readability
            paragraphs = self._split_into_paragraphs(content_text)
            
            for i, para_text in enumerate(paragraphs[:3]):  # Limit to 3 paragraphs
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = para_text
                p.font.size = Pt(18)
                p.font.color.rgb = self.colors['text']
                p.alignment = PP_ALIGN.LEFT
                p.line_spacing = 1.3
                p.space_after = Pt(12)  # Space between paragraphs
                
                # First paragraph slightly larger
                if i == 0:
                    p.font.size = Pt(20)
                    p.font.bold = True
    
    def _split_into_paragraphs(self, text: str) -> list:
        """Split long text into readable paragraphs"""
        # Split by common paragraph indicators
        sentences = text.replace('. ', '.|').replace('? ', '?|').replace('! ', '!|').split('|')
        paragraphs = []
        current_para = ""
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
                
            # If adding this sentence would make paragraph too long, start new one
            if len(current_para) + len(sentence) > 300 and current_para:
                paragraphs.append(current_para.strip())
                current_para = sentence
            else:
                current_para += (" " + sentence) if current_para else sentence
        
        # Add final paragraph
        if current_para.strip():
            paragraphs.append(current_para.strip())
        
        return paragraphs
    
    def _format_bullet_point(self, paragraph, level=0):
        """Format individual bullet point with adaptive sizing and proper spacing"""
        # Enhanced base font sizes for different levels
        base_font_sizes = [18, 16, 14, 12]  # More granular sizing
        colors = [self.colors['text'], self.colors['light_text'], self.colors['light_text'], self.colors['light_text']]
        
        # Adjust font size based on text length for better fitting
        text_length = len(paragraph.text)
        
        # Calculate adaptive font size with improved scaling
        level_idx = min(level, len(base_font_sizes) - 1)
        base_size = base_font_sizes[level_idx]
        
        # More sophisticated adaptive sizing
        if text_length > 200:
            font_size = max(base_size - 6, 10)  # Very long text
        elif text_length > 150:
            font_size = max(base_size - 4, 12)  # Long text
        elif text_length > 100:
            font_size = max(base_size - 2, 14)  # Medium text
        elif text_length > 50:
            font_size = max(base_size - 1, 16)  # Short-medium text
        else:
            font_size = base_size  # Short text
            
        # Enhanced paragraph formatting
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = colors[min(level, len(colors) - 1)]
        
        # Improved spacing and formatting
        paragraph.space_before = Pt(3)  # Space before paragraph
        paragraph.space_after = Pt(6) if level == 0 else Pt(3)  # More space after main points
        paragraph.line_spacing = 1.2  # Better line spacing
        
        # Enhanced text formatting based on level
        if level == 0:
            paragraph.font.bold = True
        elif level == 1:
            paragraph.font.italic = True  # Italics for sub-points
        
        # Add proper indentation
        paragraph.left_indent = Inches(0.25 * level)  # Progressive indentation
    
    def _calculate_text_density(self, slide_data: Slide) -> float:
        """Calculate text density ratio for adaptive layout decisions"""
        total_length = len(slide_data.title) if slide_data.title else 0
        content_items = 0
        
        if slide_data.content:
            content_items = len(slide_data.content)
            for point in slide_data.content:
                if hasattr(point, 'text'):
                    if hasattr(point.text, '__len__'):
                        total_length += len(point.text)
                    else:
                        total_length += len(str(point.text))
                    if hasattr(point, 'sub_points') and point.sub_points:
                        content_items += len(point.sub_points)
                        for sub in point.sub_points:
                            total_length += len(str(sub))
                else:
                    total_length += len(str(point))
        elif slide_data.bullets:
            content_items = len(slide_data.bullets)
            total_length += sum(len(str(bullet)) for bullet in slide_data.bullets)
        
        # Calculate density ratio (0-1 scale)
        # Base calculation on both text length and number of items
        length_factor = min(total_length / 1000, 1.0)  # Normalize to 1000 chars
        items_factor = min(content_items / 10, 1.0)    # Normalize to 10 items
        
        # Weighted average (text length weighted more heavily)
        density = (length_factor * 0.7) + (items_factor * 0.3)
        return min(density, 1.0)
    
    def _is_content_dense(self, slide_data: Slide) -> bool:
        """Determine if slide content is very dense and needs special handling"""
        density = self._calculate_text_density(slide_data)
        return density > 0.6  # Threshold for dense content
    
    def _create_overflow_safe_textbox(self, slide, slide_data: Slide):
        """Create a custom text box that ensures content never overflows slide boundaries"""
        from pptx.util import Pt, Inches
        
        # Create a custom text box with strict boundaries
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        
        logger.info(f"Creating overflow-safe text box for dense content")
        
        # Add content with tight formatting
        if slide_data.content:
            # Use more aggressive truncation for dense slides
            max_points = min(6, len(slide_data.content))
            for i, point in enumerate(slide_data.content[:max_points]):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                # Extract point text with strict length limits
                if hasattr(point, 'text'):
                    if hasattr(point.text, '__len__'):
                        point_text = point.text[:150] + "..." if len(point.text) > 150 else point.text
                    else:
                        point_text = str(point.text)
                else:
                    point_text = str(point)[:150] + "..." if len(str(point)) > 150 else str(point)
                    
                p.text = point_text
                p.level = 0
                
                # Apply compact formatting
                font = p.font
                font.size = Pt(16)  # Smaller font for dense slides
                font.bold = True
                
                # Handle sub-points more aggressively for dense slides
                if hasattr(point, 'sub_points') and point.sub_points:
                    # Limit to 2 sub-points for dense slides
                    max_sub = min(2, len(point.sub_points))
                    for j, sub in enumerate(point.sub_points[:max_sub]):
                        sub_p = tf.add_paragraph()
                        sub_text = str(sub)[:100] + "..." if len(str(sub)) > 100 else str(sub)
                        sub_p.text = sub_text
                        sub_p.level = 1
                        
                        # Apply compact formatting for sub-points
                        sub_font = sub_p.font
                        sub_font.size = Pt(14)
                        sub_font.bold = False
            
            # Add ellipsis if content was truncated
            if len(slide_data.content) > max_points:
                p = tf.add_paragraph()
                p.text = "(Additional content has been condensed)"
                p.font.italic = True
                p.font.size = Pt(12)
        
        elif slide_data.bullets:
            # Process simple bullets with strict truncation
            max_bullets = min(8, len(slide_data.bullets))
            for i, bullet in enumerate(slide_data.bullets[:max_bullets]):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                # Truncate long bullets
                bullet_text = str(bullet)[:150] + "..." if len(str(bullet)) > 150 else str(bullet)
                p.text = bullet_text
                
                # Apply compact formatting
                font = p.font
                font.size = Pt(16)
            
            # Add ellipsis if bullets were truncated
            if len(slide_data.bullets) > max_bullets:
                p = tf.add_paragraph()
                p.text = "(Additional bullets have been condensed)"
                p.font.italic = True
                p.font.size = Pt(12)
    
    def _add_image_to_slide(self, slide, image_url, position='right', custom_pos=None):
        """Add image to slide with enhanced positioning and adaptive sizing"""
        try:
            # Download image
            response = requests.get(image_url, timeout=10)
            response.raise_for_status()
            
            image_stream = BytesIO(response.content)
            
            # Calculate position with custom positioning support
            if custom_pos:
                left, top, width, height = custom_pos
            else:
                # Default position calculations with improved spacing
                if position == 'right':
                    left = Inches(5.6)  # Slightly more right for better spacing
                    top = Inches(1.6)   # Slightly higher
                    width = Inches(3.9)
                    height = Inches(4.8)
                elif position == 'center':
                    left = Inches(2.8)
                    top = Inches(4.2)   # Below text content with more margin
                    width = Inches(4.4)
                    height = Inches(2.8)
                else:  # left
                    left = Inches(0.2)
                    top = Inches(1.6)
                    width = Inches(3.9)
                    height = Inches(4.8)
            
            # Validate image dimensions to prevent off-slide placement
            slide_width = Inches(10)
            slide_height = Inches(7.5)
            
            if left + width > slide_width:
                width = slide_width - left - Inches(0.1)  # Leave small margin
            if top + height > slide_height:
                height = slide_height - top - Inches(0.1)  # Leave small margin
            
            # Add image with validated dimensions
            slide.shapes.add_picture(image_stream, left, top, width, height)
            
        except Exception as e:
            logger.error(f"Failed to add image: {e}")
            # Add placeholder with same positioning logic
            if custom_pos:
                self._add_image_placeholder(slide, position, custom_pos)
            else:
                self._add_image_placeholder(slide, position)
    
    def _add_compact_image_to_slide(self, slide, image_url, custom_pos=None):
        """Add smaller image for mixed layouts with enhanced positioning"""
        try:
            response = requests.get(image_url, timeout=10)
            response.raise_for_status()
            
            image_stream = BytesIO(response.content)
            
            # Use custom positioning if provided, otherwise default
            if custom_pos:
                left, top, width, height = custom_pos
            else:
                left = Inches(6.5)  # Adjusted for better spacing
                top = Inches(1.6)
                width = Inches(2.8)
                height = Inches(2.0)
            
            # Validate dimensions to prevent overflow
            slide_width = Inches(10)
            slide_height = Inches(7.5)
            
            if left + width > slide_width:
                width = slide_width - left - Inches(0.1)
            if top + height > slide_height:
                height = slide_height - top - Inches(0.1)
            
            slide.shapes.add_picture(image_stream, left, top, width, height)
            
        except Exception as e:
            logger.error(f"Failed to add compact image: {e}")
            if custom_pos:
                self._add_compact_image_placeholder(slide, custom_pos)
            else:
                self._add_compact_image_placeholder(slide)
    
    def _add_compact_image_placeholder(self, slide, custom_pos=None):
        """Add compact image placeholder with enhanced styling"""
        if custom_pos:
            left, top, width, height = custom_pos
        else:
            left, top, width, height = Inches(6.5), Inches(1.6), Inches(2.8), Inches(2.0)
        
        # Validate dimensions
        slide_width = Inches(10)
        slide_height = Inches(7.5)
        
        if left + width > slide_width:
            width = slide_width - left - Inches(0.1)
        if top + height > slide_height:
            height = slide_height - top - Inches(0.1)
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(248, 249, 250)
        shape.line.color.rgb = RGBColor(206, 212, 218)
        shape.line.width = Pt(1)
        
        text_frame = shape.text_frame
        text_frame.text = "🖼️\nImage"
        text_frame.margin_left = Pt(8)
        text_frame.margin_right = Pt(8)
        text_frame.margin_top = Pt(8)
        text_frame.margin_bottom = Pt(8)
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(10)
            paragraph.font.color.rgb = RGBColor(108, 117, 125)
        
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    def _add_image_placeholder(self, slide, position='right', custom_pos=None):
        """Add enhanced image placeholder when image fails to load"""
        if custom_pos:
            left, top, width, height = custom_pos
        else:
            if position == 'right':
                left, top, width, height = Inches(5.6), Inches(1.6), Inches(3.9), Inches(4.8)
            elif position == 'center':
                left, top, width, height = Inches(2.8), Inches(4.2), Inches(4.4), Inches(2.8)
            else:  # left
                left, top, width, height = Inches(0.2), Inches(1.6), Inches(3.9), Inches(4.8)
        
        # Validate placeholder dimensions
        slide_width = Inches(10)
        slide_height = Inches(7.5)
        
        if left + width > slide_width:
            width = slide_width - left - Inches(0.1)
        if top + height > slide_height:
            height = slide_height - top - Inches(0.1)
        
        # Create enhanced placeholder
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(248, 249, 250)  # Light gray
        shape.line.color.rgb = RGBColor(206, 212, 218)  # Darker border
        shape.line.width = Pt(1.5)
        
        # Add enhanced placeholder text with icon-like appearance
        text_frame = shape.text_frame
        text_frame.text = "🖼️\nImage Placeholder"
        text_frame.margin_left = Pt(12)
        text_frame.margin_right = Pt(12)
        text_frame.margin_top = Pt(12)
        text_frame.margin_bottom = Pt(12)
        
        # Format text
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(14) if width > Inches(3) else Pt(12)
            paragraph.font.color.rgb = RGBColor(108, 117, 125)  # Medium gray
        
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.text = "Image Placeholder"
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = RGBColor(100, 100, 100)
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    def _add_diagram_to_slide(self, slide, slide_data: Slide):
        """Add diagram to slide"""
        diagram_type = slide_data.diagram_type
        diagram_data = slide_data.diagram_data or []
        
        if diagram_type == 'process':
            self._create_process_diagram(slide, diagram_data)
        elif diagram_type == 'comparison':
            self._create_comparison_diagram(slide, diagram_data)
        elif diagram_type == 'hierarchy':
            self._create_hierarchy_diagram(slide, diagram_data)
    
    def _create_process_diagram(self, slide, steps):
        """Create enhanced process flow diagram with adaptive sizing and positioning"""
        if not steps:
            return
        
        # Enhanced diagram dimensions with better spacing
        max_steps_per_row = min(len(steps), 4)  # Limit for better visibility
        
        # Adaptive sizing based on number of steps and content density
        if max_steps_per_row <= 2:
            step_width = Inches(2.2)
            step_height = Inches(0.9)
            arrow_width = Inches(0.6)
        elif max_steps_per_row == 3:
            step_width = Inches(2.0)
            step_height = Inches(0.8)
            arrow_width = Inches(0.5)
        else:  # 4 steps
            step_width = Inches(1.7)
            step_height = Inches(0.7)
            arrow_width = Inches(0.4)
        
        # Calculate total width with improved spacing
        total_width = max_steps_per_row * step_width + (max_steps_per_row - 1) * arrow_width
        
        # Ensure diagram fits within slide boundaries
        slide_content_width = Inches(9.2)  # Leave margins
        if total_width > slide_content_width:
            # Proportionally reduce dimensions
            scale_factor = slide_content_width / total_width
            step_width = step_width * scale_factor
            arrow_width = arrow_width * scale_factor
            total_width = slide_content_width
        
        # Center the diagram horizontally with proper margins
        start_left = (Inches(10) - total_width) / 2
        top = Inches(4.3)  # Position with adequate spacing below text
        
        for i, step in enumerate(steps[:max_steps_per_row]):
            # Calculate position
            left = start_left + i * (step_width + arrow_width)
            
            # Create step box with enhanced styling
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, step_width, step_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.colors['secondary']
            shape.line.color.rgb = self.colors['primary']
            shape.line.width = Pt(2)
            
            # Add rounded corners effect with shadow-like border
            shape.shadow.inherit = False
            shape.shadow.style = 1  # Outer shadow
            shape.shadow.blur_radius = Pt(3)
            shape.shadow.distance = Pt(2)
            shape.shadow.color.rgb = RGBColor(0, 0, 0)
            shape.shadow.transparency = 0.3
            
            # Enhanced text formatting
            text_frame = shape.text_frame
            text_frame.margin_left = Pt(8)
            text_frame.margin_right = Pt(8)
            text_frame.margin_top = Pt(6)
            text_frame.margin_bottom = Pt(6)
            text_frame.word_wrap = True
            
            # Adaptive text handling
            step_text = step.get('step', str(step)) if isinstance(step, dict) else str(step)
            
            # Intelligent text truncation based on step box size
            max_chars = int(step_width.inches * 12)  # Approximate chars per inch
            if len(step_text) > max_chars:
                step_text = step_text[:max_chars-3] + "..."
            
            text_frame.text = step_text
            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(11) if step_width >= Inches(1.8) else Pt(10)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.bold = True
            paragraph.line_spacing = 1.1
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Add enhanced arrow between steps
            if i < min(len(steps), max_steps_per_row) - 1:
                arrow_left = left + step_width
                arrow_top = top + (step_height - Inches(0.3)) / 2  # Center vertically
                
                arrow_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW, arrow_left, arrow_top, 
                    arrow_width, Inches(0.3)
                )
                arrow_shape.fill.solid()
                arrow_shape.fill.fore_color.rgb = self.colors['accent']
                arrow_shape.line.color.rgb = self.colors['primary']
                arrow_shape.line.width = Pt(1)
        
        # Add continuation indicator if there are more steps
        if len(steps) > max_steps_per_row:
            remaining_steps = len(steps) - max_steps_per_row
            indicator_left = start_left + max_steps_per_row * (step_width + arrow_width) - arrow_width
            indicator_shape = slide.shapes.add_textbox(
                indicator_left, top + step_height + Inches(0.1), Inches(1), Inches(0.3)
            )
            indicator_frame = indicator_shape.text_frame
            indicator_frame.text = f"... +{remaining_steps} more"
            indicator_frame.paragraphs[0].font.size = Pt(10)
            indicator_frame.paragraphs[0].font.italic = True
            indicator_frame.paragraphs[0].font.color.rgb = self.colors['light_text']
            indicator_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _add_compact_diagram_to_slide(self, slide, slide_data: Slide, offset_top=None):
        """Add compact diagram for mixed layouts with positioning control"""
        diagram_type = slide_data.diagram_type
        diagram_data = slide_data.diagram_data or []
        
        # Calculate diagram position with offset support
        base_top = Inches(4.2)
        if offset_top:
            diagram_top = Inches(1.4) + offset_top
        else:
            diagram_top = base_top
        
        if diagram_type == 'process':
            self._create_compact_process_diagram(slide, diagram_data, diagram_top)
        elif diagram_type == 'comparison':
            self._create_compact_comparison_diagram(slide, diagram_data, diagram_top)
    
    def _create_compact_process_diagram(self, slide, steps, top=None):
        """Create compact process diagram for mixed layouts with positioning control"""
        if not steps:
            return
        
        max_steps = min(len(steps), 3)  # Max 3 steps for compact version
        step_width = Inches(1.3)
        step_height = Inches(0.6)
        arrow_width = Inches(0.25)
        
        # Calculate total width and center positioning
        total_width = max_steps * step_width + (max_steps - 1) * arrow_width
        start_left = (Inches(10) - total_width) / 2
        diagram_top = top if top else Inches(4.2)
        
        for i, step in enumerate(steps[:max_steps]):
            left = start_left + i * (step_width + arrow_width)
            
            # Create step box with enhanced compact styling
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, diagram_top, step_width, step_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.colors['secondary']
            shape.line.color.rgb = self.colors['primary']
            shape.line.width = Pt(1.5)
            
            # Enhanced text formatting for compact diagrams
            text_frame = shape.text_frame
            text_frame.margin_left = Pt(4)
            text_frame.margin_right = Pt(4)
            text_frame.margin_top = Pt(3)
            text_frame.margin_bottom = Pt(3)
            text_frame.word_wrap = True
            
            step_text = step.get('step', str(step)) if isinstance(step, dict) else str(step)
            if len(step_text) > 12:
                step_text = step_text[:10] + ".."
            
            text_frame.text = step_text
            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(9)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.bold = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Add arrow between steps
            if i < max_steps - 1:
                arrow_left = left + step_width
                arrow_top = diagram_top + (step_height - Inches(0.2)) / 2
                arrow_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW, arrow_left, arrow_top, 
                    arrow_width, Inches(0.2)
                )
                arrow_shape.fill.solid()
                arrow_shape.fill.fore_color.rgb = self.colors['accent']
                arrow_shape.line.color.rgb = self.colors['accent']
        
        for i, step in enumerate(steps[:max_steps]):
            left = start_left + i * (step_width + arrow_width)
            
            # Add step box
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, step_width, step_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.colors['secondary']
            shape.line.color.rgb = self.colors['primary']
            
            # Add step text
            text_frame = shape.text_frame
            step_text = step.get('step', str(step)) if isinstance(step, dict) else str(step)
            if len(step_text) > 10:
                step_text = step_text[:8] + ".."
            
            text_frame.text = step_text
            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(8)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.bold = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Add arrow
            if i < max_steps - 1:
                arrow_left = left + step_width
                arrow_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW, arrow_left, top + Inches(0.15), 
                    arrow_width, Inches(0.2)
                )
                arrow_shape.fill.solid()
                arrow_shape.fill.fore_color.rgb = self.colors['accent']
    
    def _create_comparison_diagram(self, slide, comparison_data):
        """Create enhanced comparison diagram with adaptive sizing and styling"""
        if len(comparison_data) < 2:
            return
        
        # Enhanced box dimensions with better proportions
        box_width = Inches(3.8)
        box_height = Inches(2.8)
        gap = Inches(0.6)
        
        # Center the comparison boxes with proper margins
        total_width = 2 * box_width + gap
        start_left = (Inches(10) - total_width) / 2
        top = Inches(4.2)  # Positioned to avoid text overlap
        
        # Left column with enhanced styling
        left_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, start_left, top, box_width, box_height
        )
        left_box.fill.solid()
        left_box.fill.fore_color.rgb = RGBColor(230, 242, 255)  # Light blue
        left_box.line.color.rgb = self.colors['secondary']
        left_box.line.width = Pt(2.5)
        
        # Add subtle shadow effect
        left_box.shadow.inherit = False
        left_box.shadow.style = 1
        left_box.shadow.blur_radius = Pt(4)
        left_box.shadow.distance = Pt(3)
        left_box.shadow.color.rgb = RGBColor(0, 0, 0)
        left_box.shadow.transparency = 0.25
        
        # Enhanced text formatting for left box
        left_text = left_box.text_frame
        left_text.margin_left = Pt(15)
        left_text.margin_right = Pt(15)
        left_text.margin_top = Pt(15)
        left_text.margin_bottom = Pt(15)
        left_text.word_wrap = True
        
        # Adaptive text handling
        left_title = comparison_data[0].get('title', 'Option A') if isinstance(comparison_data[0], dict) else str(comparison_data[0])
        if len(left_title) > 60:
            left_title = left_title[:57] + "..."
        
        left_text.text = left_title
        self._format_comparison_text(left_text)
        
        # Right column with complementary styling
        right_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, start_left + box_width + gap, top, box_width, box_height
        )
        right_box.fill.solid()
        right_box.fill.fore_color.rgb = RGBColor(255, 242, 230)  # Light orange
        right_box.line.color.rgb = self.colors['accent']
        right_box.line.width = Pt(2.5)
        
        # Add shadow effect to right box
        right_box.shadow.inherit = False
        right_box.shadow.style = 1
        right_box.shadow.blur_radius = Pt(4)
        right_box.shadow.distance = Pt(3)
        right_box.shadow.color.rgb = RGBColor(0, 0, 0)
        right_box.shadow.transparency = 0.25
        
        # Enhanced text formatting for right box
        right_text = right_box.text_frame
        right_text.margin_left = Pt(15)
        right_text.margin_right = Pt(15)
        right_text.margin_top = Pt(15)
        right_text.margin_bottom = Pt(15)
        right_text.word_wrap = True
        
        right_title = comparison_data[1].get('title', 'Option B') if isinstance(comparison_data[1], dict) else str(comparison_data[1])
        if len(right_title) > 60:
            right_title = right_title[:57] + "..."
        
        right_text.text = right_title
        self._format_comparison_text(right_text)
        
        # Add connecting element (vs. indicator)
        vs_left = start_left + box_width + (gap / 2) - Inches(0.25)
        vs_top = top + (box_height / 2) - Inches(0.15)
        vs_shape = slide.shapes.add_textbox(vs_left, vs_top, Inches(0.5), Inches(0.3))
        vs_frame = vs_shape.text_frame
        vs_frame.text = "VS"
        vs_paragraph = vs_frame.paragraphs[0]
        vs_paragraph.alignment = PP_ALIGN.CENTER
        vs_paragraph.font.size = Pt(12)
        vs_paragraph.font.bold = True
        vs_paragraph.font.color.rgb = self.colors['primary']
        vs_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    def _create_compact_comparison_diagram(self, slide, comparison_data, top=None):
        """Create compact comparison diagram for mixed layouts with positioning control"""
        if len(comparison_data) < 2:
            return
        
        # Enhanced compact dimensions
        box_width = Inches(2.8)
        box_height = Inches(1.6)
        gap = Inches(0.4)
        
        # Center the boxes horizontally
        total_width = 2 * box_width + gap
        start_left = (Inches(10) - total_width) / 2
        diagram_top = top if top else Inches(4.5)
        
        # Left column with enhanced styling
        left_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, start_left, diagram_top, box_width, box_height
        )
        left_box.fill.solid()
        left_box.fill.fore_color.rgb = RGBColor(230, 242, 255)
        left_box.line.color.rgb = self.colors['secondary']
        left_box.line.width = Pt(1.5)
        
        left_text = left_box.text_frame
        left_text.margin_left = Pt(8)
        left_text.margin_right = Pt(8)
        left_text.margin_top = Pt(6)
        left_text.margin_bottom = Pt(6)
        left_text.word_wrap = True
        
        left_title = comparison_data[0].get('title', 'Option A') if isinstance(comparison_data[0], dict) else str(comparison_data[0])
        if len(left_title) > 35:
            left_title = left_title[:32] + "..."
        left_text.text = left_title
        
        paragraph = left_text.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(11)
        paragraph.font.bold = True
        paragraph.font.color.rgb = self.colors['primary']
        paragraph.line_spacing = 1.1
        left_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Right column with complementary styling
        right_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, start_left + box_width + gap, diagram_top, box_width, box_height
        )
        right_box.fill.solid()
        right_box.fill.fore_color.rgb = RGBColor(255, 242, 230)
        right_box.line.color.rgb = self.colors['accent']
        right_box.line.width = Pt(1.5)
        
        right_text = right_box.text_frame
        right_text.margin_left = Pt(8)
        right_text.margin_right = Pt(8)
        right_text.margin_top = Pt(6)
        right_text.margin_bottom = Pt(6)
        right_text.word_wrap = True
        
        right_title = comparison_data[1].get('title', 'Option B') if isinstance(comparison_data[1], dict) else str(comparison_data[1])
        if len(right_title) > 35:
            right_title = right_title[:32] + "..."
        right_text.text = right_title
        
        paragraph = right_text.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(11)
        paragraph.font.bold = True
        paragraph.font.color.rgb = self.colors['primary']
        paragraph.line_spacing = 1.1
        right_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        right_box.fill.solid()
        right_box.fill.fore_color.rgb = RGBColor(250, 240, 230)
        right_box.line.color.rgb = self.colors['accent']
        
        right_text = right_box.text_frame
        right_title = comparison_data[1].get('title', 'Option B') if isinstance(comparison_data[1], dict) else str(comparison_data[1])
        if len(right_title) > 30:
            right_title = right_title[:27] + "..."
        right_text.text = right_title
        
        paragraph = right_text.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(12)
        paragraph.font.bold = True
        paragraph.font.color.rgb = self.colors['primary']
        right_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    def _create_hierarchy_diagram(self, slide, hierarchy_data):
        """Create hierarchy diagram with proper sizing"""
        if not hierarchy_data:
            return
        
        # Top level box
        top_item = hierarchy_data[0] if hierarchy_data else {}
        top_text = top_item.get('title', 'Root') if isinstance(top_item, dict) else str(top_item)
        
        # Ensure text fits in box
        if len(top_text) > 20:
            top_text = top_text[:17] + "..."
        
        # Center the top box
        box_width = Inches(3)
        box_height = Inches(0.8)
        top_left = (Inches(10) - box_width) / 2
        top_top = Inches(4)
        
        top_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, top_left, top_top, box_width, box_height
        )
        top_box.fill.solid()
        top_box.fill.fore_color.rgb = self.colors['primary']
        top_box.line.color.rgb = self.colors['primary']
        top_box.line.width = Pt(2)
        
        text_frame = top_box.text_frame
        text_frame.margin_left = Pt(8)
        text_frame.margin_right = Pt(8)
        text_frame.margin_top = Pt(8)
        text_frame.margin_bottom = Pt(8)
        
        text_frame.text = top_text
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.font.bold = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Add sub-items if available
        if len(hierarchy_data) > 1:
            sub_items = hierarchy_data[1:4]  # Max 3 sub-items
            sub_box_width = Inches(2.2)
            sub_box_height = Inches(0.6)
            sub_gap = Inches(0.4)
            
            total_sub_width = len(sub_items) * sub_box_width + (len(sub_items) - 1) * sub_gap
            sub_start_left = (Inches(10) - total_sub_width) / 2
            sub_top = top_top + box_height + Inches(0.5)
            
            for i, sub_item in enumerate(sub_items):
                sub_text = sub_item.get('title', f'Item {i+1}') if isinstance(sub_item, dict) else str(sub_item)
                if len(sub_text) > 15:
                    sub_text = sub_text[:12] + "..."
                
                sub_left = sub_start_left + i * (sub_box_width + sub_gap)
                
                sub_box = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, sub_left, sub_top, sub_box_width, sub_box_height
                )
                sub_box.fill.solid()
                sub_box.fill.fore_color.rgb = self.colors['secondary']
                sub_box.line.color.rgb = self.colors['primary']
                
                sub_text_frame = sub_box.text_frame
                sub_text_frame.text = sub_text
                sub_paragraph = sub_text_frame.paragraphs[0]
                sub_paragraph.alignment = PP_ALIGN.CENTER
                sub_paragraph.font.size = Pt(11)
                sub_paragraph.font.color.rgb = RGBColor(255, 255, 255)
                sub_paragraph.font.bold = True
                sub_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                # Add connecting line
                line_start_x = top_left + box_width / 2
                line_start_y = top_top + box_height
                line_end_x = sub_left + sub_box_width / 2
                line_end_y = sub_top
                
                connector = slide.shapes.add_connector(
                    1, line_start_x, line_start_y, line_end_x, line_end_y
                )
                connector.line.color.rgb = self.colors['primary']
                connector.line.width = Pt(2)
    
    def _format_comparison_text(self, text_frame):
        """Format comparison diagram text with enhanced styling"""
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(15)
        paragraph.font.bold = True
        paragraph.font.color.rgb = self.colors['primary']
        paragraph.line_spacing = 1.2
        paragraph.space_before = Pt(3)
        paragraph.space_after = Pt(3)
        
        # Enhanced text frame properties
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    def _add_legacy_diagram(self, slide, diagram_description):
        """Add legacy diagram support"""
        try:
            diagram_img = generate_diagram_image(diagram_description)
            slide.shapes.add_picture(diagram_img, Inches(2), Inches(2), width=Inches(4))
        except Exception as e:
            logger.error(f"Failed to embed legacy diagram: {diagram_description} - {e}")
